from __future__ import annotations

import argparse
import json
import shutil
import tempfile
from copy import deepcopy
from pathlib import Path
from typing import Any

import fitz
from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = ROOT / "assets" / "shunyejiang-template.pptx"

TEXT_SHAPES = {
    "page": "TextBox 2",
    "total": "TextBox 3",
    "status": "TextBox 5",
    "keyword_1": "TextBox 9",
    "keyword_2": "TextBox 11",
    "keyword_3": "TextBox 13",
    "ipad_page": "TextBox 14",
    "screen_page": "TextBox 15",
    "script": "TextBox 18",
    "next_step": "TextBox 20",
}
THUMBNAIL_SHAPE = "Picture 7"
CONTENT_PANEL_SHAPE = "Rectangle 17"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Render a page-synced companion PPT from a source PDF and page-plan JSON.")
    parser.add_argument("--source-pdf", required=True, type=Path)
    parser.add_argument("--plan-json", required=True, type=Path)
    parser.add_argument("--output-pptx", required=True, type=Path)
    parser.add_argument("--template-pptx", default=DEFAULT_TEMPLATE, type=Path)
    parser.add_argument("--total-pages", type=int, default=None, help="Displayed total page count. Defaults to source PDF page count.")
    return parser.parse_args()


def load_plan(path: Path) -> list[dict[str, Any]]:
    data = json.loads(path.read_text(encoding="utf-8"))
    pages = data.get("pages", data) if isinstance(data, dict) else data
    if not isinstance(pages, list) or not pages:
        raise ValueError("plan JSON must contain a non-empty pages list")
    normalized: list[dict[str, Any]] = []
    for index, item in enumerate(pages, start=1):
        if not isinstance(item, dict):
            raise ValueError(f"page item {index} must be an object")
        page = item.get("page", index)
        if not isinstance(page, int) or page < 1:
            raise ValueError(f"page item {index}: page must be a positive integer")
        keywords = item.get("keywords", [])
        if isinstance(keywords, str):
            keywords = [part.strip() for part in keywords.replace("，", ",").split(",") if part.strip()]
        keywords = [str(k).strip() for k in keywords if str(k).strip()][:3]
        while len(keywords) < 3:
            keywords.append("")
        normalized.append(
            {
                "page": page,
                "status": str(item.get("status", "")).strip(),
                "keywords": keywords,
                "script": str(item.get("script", "")).strip(),
                "next_step": str(item.get("next_step", "")).strip(),
            }
        )
    return normalized


def find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"template shape not found: {name}")


def set_text_preserve_style(shape, text: str) -> None:
    """Replace text while preserving the template's first-run style."""
    if not hasattr(shape, "text_frame"):
        return
    frame = shape.text_frame
    first_paragraph = frame.paragraphs[0]
    if first_paragraph.runs:
        first_run = first_paragraph.runs[0]
    else:
        first_run = first_paragraph.add_run()
    first_run.text = text
    for run in first_paragraph.runs[1:]:
        run.text = ""
    for paragraph in frame.paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def center_script_box(slide) -> None:
    """Keep the main read-aloud text visually centered in the right content panel."""
    panel = find_shape(slide, CONTENT_PANEL_SHAPE)
    script = find_shape(slide, TEXT_SHAPES["script"])
    script.left = int(panel.left + (panel.width - script.width) / 2)
    script.top = int(panel.top + (panel.height - script.height) / 2)
    if hasattr(script, "text_frame"):
        script.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE


def duplicate_template_slide(prs: Presentation, source_slide, skip_picture_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    copy_slide_background(source_slide, slide)
    for shape in source_slide.shapes:
        if shape.name == skip_picture_name:
            continue
        slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return slide


def copy_slide_background(source_slide, target_slide) -> None:
    source_background = source_slide._element.cSld.find(qn("p:bg"))
    target_background = target_slide._element.cSld.find(qn("p:bg"))
    if target_background is not None:
        target_slide._element.cSld.remove(target_background)
    if source_background is not None:
        target_slide._element.cSld.insert(0, deepcopy(source_background))


def delete_slide(prs: Presentation, slide) -> None:
    slide_id = slide.slide_id
    slide_id_list = prs.slides._sldIdLst
    slide_id_element = slide_id_list.xpath(f'./p:sldId[@id="{slide_id}"]')[0]
    prs.part.drop_rel(slide_id_element.rId)
    slide_id_list.remove(slide_id_element)


def make_thumbnail(source_path: Path, output_path: Path, width: int, height: int) -> None:
    with Image.open(source_path) as img:
        img = img.convert("RGB")
        img.thumbnail((width, height), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (width, height), (31, 32, 29))
        canvas.paste(img, ((width - img.width) // 2, (height - img.height) // 2))
        canvas.save(output_path, quality=90)


def render_pdf_pages(source_pdf: Path, page_numbers: list[int], tmp_dir: Path, thumb_px: tuple[int, int]) -> tuple[dict[int, Path], int]:
    doc = fitz.open(str(source_pdf))
    page_count = doc.page_count
    rendered: dict[int, Path] = {}
    for page_number in sorted(set(page_numbers)):
        if page_number > page_count:
            raise ValueError(f"page {page_number} exceeds source page count {page_count}")
        page = doc[page_number - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(0.8, 0.8), alpha=False)
        raw_path = tmp_dir / f"raw_{page_number:03d}.jpg"
        pix.save(str(raw_path))
        thumb_path = tmp_dir / f"thumb_{page_number:03d}.jpg"
        make_thumbnail(raw_path, thumb_path, *thumb_px)
        rendered[page_number] = thumb_path
    return rendered, page_count


def render(source_pdf: Path, plan_json: Path, template_pptx: Path, output_pptx: Path, displayed_total: int | None = None) -> None:
    if not source_pdf.exists():
        raise FileNotFoundError(source_pdf)
    if not template_pptx.exists():
        raise FileNotFoundError(template_pptx)

    page_plan = load_plan(plan_json)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    tmp_dir = Path(tempfile.mkdtemp(prefix="companion_ppt_"))
    try:
        prs = Presentation(str(template_pptx))
        template_slide = prs.slides[0]
        picture_shape = find_shape(template_slide, THUMBNAIL_SHAPE)
        picture_box = {
            "left": picture_shape.left,
            "top": picture_shape.top,
            "width": picture_shape.width,
            "height": picture_shape.height,
        }
        thumb_px = (max(1, int(picture_shape.width / 914400 * 360)), max(1, int(picture_shape.height / 914400 * 360)))
        thumbnails, source_total = render_pdf_pages(source_pdf, [item["page"] for item in page_plan], tmp_dir, thumb_px)
        total = displayed_total or source_total

        slides = [duplicate_template_slide(prs, template_slide, THUMBNAIL_SHAPE) for _ in page_plan]
        delete_slide(prs, template_slide)

        for slide, item in zip(slides, page_plan):
            center_script_box(slide)
            page_number = item["page"]
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["page"]), f"{page_number:02d}")
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["total"]), f"/ {total}")
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["status"]), item["status"])
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["keyword_1"]), item["keywords"][0])
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["keyword_2"]), item["keywords"][1])
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["keyword_3"]), item["keywords"][2])
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["ipad_page"]), f"iPad 第 {page_number} 页")
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["screen_page"]), f"对应大屏第 {page_number} 页")
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["script"]), item["script"])
            set_text_preserve_style(find_shape(slide, TEXT_SHAPES["next_step"]), item["next_step"])
            slide.shapes.add_picture(
                str(thumbnails[page_number]),
                picture_box["left"],
                picture_box["top"],
                picture_box["width"],
                picture_box["height"],
            )

        prs.save(str(output_pptx))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def main() -> None:
    args = parse_args()
    render(args.source_pdf, args.plan_json, args.template_pptx, args.output_pptx, args.total_pages)
    print(args.output_pptx)


if __name__ == "__main__":
    main()
